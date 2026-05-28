#!/usr/bin/env python
"""Build defense PPT based on visual design spec."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

# === COLORS ===
DARK_BG = RGBColor(0x0D, 0x1B, 0x2A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT_BLUE = RGBColor(0x29, 0x80, 0xB9)
TEXT_DARK = RGBColor(0x1A, 0x2A, 0x3A)
TEXT_GRAY = RGBColor(0x56, 0x65, 0x73)
LIGHT_GRAY = RGBColor(0xD5, 0xD8, 0xDC)
CARD_BG = RGBColor(0xF7, 0xF9, 0xFC)
GREEN_ACCENT = RGBColor(0x27, 0xAE, 0x60)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_BLUE_2 = RGBColor(0x34, 0x98, 0xDB)  # for decor elements
FAINT = RGBColor(0xBB, 0xBB, 0xBB)
BG_SOFT = RGBColor(0xF0, 0xF4, 0xF8)
DARKER_BG = RGBColor(0x1B, 0x3A, 0x5C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Helper functions
def add_blank_slide():
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_line(slide, left, top, width, height, color, w=Pt(1.5)):
    w_adj = max(width, Pt(2))
    h_adj = max(height, Pt(2))
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w_adj, h_adj)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_textbox(slide, left, top, width, height, text, font_size=Pt(14), color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multi_text(slide, left, top, width, height, lines, font_name='Microsoft YaHei', anchor=MSO_ANCHOR.TOP):
    """lines = [(text, size, color, bold, alignment), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, size, color, bold, align = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
    return txBox

def add_img(slide, path, left, top, width=None, height=None):
    if width and height:
        return slide.shapes.add_picture(path, left, top, width, height)
    elif width:
        return slide.shapes.add_picture(path, left, top, width=width)
    elif height:
        return slide.shapes.add_picture(path, left, top, height=height)
    else:
        return slide.shapes.add_picture(path, left, top)

# Image paths
THESIS_IMG = 'D:/素材/Claude/portfolio/thesis-images'
CAR_MODEL = 'D:/素材/Claude/portfolio/car-model.png'
CAR_CHASSIS = 'D:/素材/Claude/portfolio/car-chassis.png'
CAR_DOCK = 'D:/素材/Claude/portfolio/car-docking.png'

# ==========================================
# SLIDE 1: COVER
# ==========================================
s1 = add_blank_slide()
set_bg(s1, DARK_BG)
# Top accent line
add_line(s1, Inches(0.8), Inches(0.9), Inches(1.5), Pt(4), ACCENT_BLUE_2)
# School
add_textbox(s1, Inches(0.8), Inches(1.1), Inches(8), Inches(0.5),
    '齐鲁工业大学 · 机械工程学院', Pt(14), RGBColor(0xFF,0xFF,0xFF), False, PP_ALIGN.LEFT, 'Microsoft YaHei')
# Title
add_textbox(s1, Inches(2.5), Inches(2.8), Inches(8.5), Inches(1.6),
    '露天矿区无人接驳小车系统设计', Pt(40), WHITE, True, PP_ALIGN.CENTER, 'Microsoft YaHei')
# Decorative line under title
add_line(s1, Inches(6.2), Inches(4.3), Inches(0.9), Pt(3), ACCENT_BLUE_2)
# Info
add_multi_text(s1, Inches(4), Inches(4.8), Inches(5.5), Inches(1.2), [
    ('答辩人：张紫然    专业班级：机器人22-1    指导教师：李老师', Pt(14), RGBColor(0xCC,0xCC,0xCC), False, PP_ALIGN.CENTER),
])
# Date
add_textbox(s1, Inches(5), Inches(6.6), Inches(3.5), Inches(0.4),
    '2026年5月', Pt(11), RGBColor(0x99,0x99,0x99), False, PP_ALIGN.CENTER)
# Bottom line
add_line(s1, Inches(0.8), Inches(7.2), Inches(11.5), Pt(1), RGBColor(0x33,0x44,0x55))

# ==========================================
# SLIDE 2: TOC
# ==========================================
s2 = add_blank_slide()
set_bg(s2, WHITE)
# Title
add_textbox(s2, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6),
    '汇报提纲', Pt(32), TEXT_DARK, True)
add_line(s2, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

toc_items = [
    ('01', '研究背景与问题', True),
    ('02', '系统总体方案', False),
    ('03', '机械结构设计', False),
    ('04', '电控与传感器系统', False),
    ('05', '运动控制与导航', False),
    ('06', '仿真验证与总结', False),
]
for i, (num, title, active) in enumerate(toc_items):
    y = Inches(1.8) + Inches(i * 0.75)
    if active:
        # Active highlight
        add_rect(s2, Inches(1.5), y, Inches(9.5), Inches(0.55), BG_SOFT, ACCENT_BLUE_2, Pt(1.5))
    else:
        add_rect(s2, Inches(1.5), y, Inches(9.5), Inches(0.55), None)
    add_textbox(s2, Inches(1.8), y + Pt(6), Inches(0.6), Inches(0.4),
        num, Pt(20), ACCENT_BLUE_2 if active else FAINT, True)
    add_textbox(s2, Inches(2.6), y + Pt(8), Inches(7), Inches(0.4),
        title, Pt(18), TEXT_DARK, True if active else False)
# Page number
add_textbox(s2, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '02', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 3: Research Background
# ==========================================
s3 = add_blank_slide()
set_bg(s3, WHITE)
add_textbox(s3, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
    '研究背景', Pt(30), TEXT_DARK, True)
add_line(s3, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

# Left column - key points
add_textbox(s3, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.5),
    '低空物流正在扩张，但末端接驳仍是薄弱环节', Pt(20), TEXT_DARK, True)
add_multi_text(s3, Inches(0.8), Inches(2.2), Inches(5.8), Inches(3.5), [
    ('· 低空经济政策推动无人机配送在物流园区、矿区等场景加速落地', Pt(16), TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', Pt(8), TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('· 无人机降落后的货物交接仍需大量人工参与', Pt(16), TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', Pt(8), TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('· 矿区环境对人员不友好——高温、粉尘、偏远', Pt(16), TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('', Pt(8), TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('· 从"无人飞行"到"无人接驳"，最后一米尚未打通', Pt(16), TEXT_GRAY, False, PP_ALIGN.LEFT),
])

# Right column - 3 keyword cards
for i, (title, desc) in enumerate([
    ('低空经济', '国家战略性\n新兴产业'),
    ('智慧物流', '无人机+\n地面机器人协同'),
    ('末端配送', '最后1公里\n→最后1米'),
]):
    y = Inches(1.5) + Inches(i * 1.8)
    add_rect(s3, Inches(7.5), y, Inches(5), Inches(1.5), CARD_BG, None)
    add_line(s3, Inches(7.5), y, Inches(5), Pt(3), ACCENT_BLUE_2)
    add_textbox(s3, Inches(7.8), y + Inches(0.2), Inches(4.4), Inches(0.5),
        title, Pt(18), TEXT_DARK, True)
    add_textbox(s3, Inches(7.8), y + Inches(0.7), Inches(4.4), Inches(0.6),
        desc, Pt(14), TEXT_GRAY, False)

add_textbox(s3, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '03', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 4: Problems
# ==========================================
s4 = add_blank_slide()
set_bg(s4, WHITE)
add_textbox(s4, Inches(0.8), Inches(0.5), Inches(8), Inches(0.6),
    '现存问题', Pt(30), TEXT_DARK, True)
add_line(s4, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)
add_textbox(s4, Inches(0.8), Inches(1.4), Inches(8), Inches(0.5),
    '三个问题，一个都没绕开', Pt(18), TEXT_GRAY, False)

problems = [
    ('降落冲击', '无人机丢货时冲击力大，\n可能损坏货物', RED_ACCENT, '安全性'),
    ('定位偏差', 'GPS偏差1-2米，\n无人机与接驳平台对不准', ACCENT_BLUE, '接驳精度'),
    ('人工对接', '最后还是需要人搬一下，\n全流程自动化断裂', GREEN_ACCENT, '效率'),
]
for i, (title, desc, color, impact) in enumerate(problems):
    x = Inches(0.8) + Inches(i * 4.2)
    add_rect(s4, x, Inches(2.5), Inches(3.8), Inches(4), CARD_BG, None)
    add_line(s4, x, Inches(2.5), Inches(3.8), Pt(4), color)
    add_textbox(s4, x + Inches(0.4), Inches(2.8), Inches(3), Inches(0.5),
        title, Pt(24), TEXT_DARK, True)
    add_textbox(s4, x + Inches(0.4), Inches(3.5), Inches(3), Inches(0.8),
        desc, Pt(15), TEXT_GRAY, False)
    add_textbox(s4, x + Inches(0.4), Inches(4.6), Inches(2), Inches(0.4),
        f'影响：{impact}', Pt(13), color, True)

add_textbox(s4, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '04', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 5: Research Objectives
# ==========================================
s5 = add_blank_slide()
set_bg(s5, WHITE)
add_textbox(s5, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '研究目标', Pt(30), TEXT_DARK, True)
add_line(s5, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

# Center statement
add_textbox(s5, Inches(2), Inches(2.2), Inches(9.5), Inches(0.6),
    '设计一台能跑、能接、能找路的无人接驳小车', Pt(23), TEXT_DARK, True, PP_ALIGN.CENTER)

# Four modules
modules = [
    ('机械系统', '四轮独立驱动底盘\n自动接驳装置\nSolidWorks三维建模', ACCENT_BLUE),
    ('控制系统', '工控机+PLC协同\n传感器选型与布局\nI/O点位与时序逻辑', GREEN_ACCENT),
    ('运动控制', '麦克纳姆轮运动学\n前馈+PID闭环调速\n精准启停控制', RED_ACCENT),
    ('导航系统', '激光SLAM建图\nA*全局路径规划\nDWA局部避障', RGBColor(0x8E,0x44,0xAD)),
]
for i, (title, desc, color) in enumerate(modules):
    x = Inches(0.8) + Inches(i * 3.1)
    add_rect(s5, x, Inches(3.6), Inches(2.8), Inches(3), CARD_BG, None)
    add_line(s5, x, Inches(3.6), Inches(2.8), Pt(4), color)
    add_textbox(s5, x + Inches(0.3), Inches(3.9), Inches(2.2), Inches(0.4),
        title, Pt(18), TEXT_DARK, True, PP_ALIGN.CENTER)
    add_textbox(s5, x + Inches(0.3), Inches(4.5), Inches(2.2), Inches(1.5),
        desc, Pt(14), TEXT_GRAY, False, PP_ALIGN.CENTER)

add_textbox(s5, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '05', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 6: Overall System Design
# ==========================================
s6 = add_blank_slide()
set_bg(s6, WHITE)
add_textbox(s6, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '系统总体方案', Pt(30), TEXT_DARK, True)
add_line(s6, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)
add_textbox(s6, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5),
    '三大系统，层层配合——机械负责"站稳接住"，控制负责"听话执行"，导航负责"找到路"', Pt(16), TEXT_GRAY, False)

# Architecture diagram - 3 layers
layers = [
    ('感知与导航层', '激光雷达 → SLAM建图 → A*全局路径 → DWA局部避障', ACCENT_BLUE_2),
    ('控制层', '工控机（上位机·跑ROS）↕ Modbus/TCP ↕ PLC（下位机·管执行）→ 电机驱动', GREEN_ACCENT),
    ('机械层', '麦克纳姆轮底盘 + 自动缓冲纠偏接驳装置', RGBColor(0xE6,0x7E,0x22)),
]
for i, (name, desc, color) in enumerate(layers):
    y = Inches(2) + Inches(i * 1.5)
    add_rect(s6, Inches(1.5), y, Inches(10.5), Inches(1.2), CARD_BG, color, Pt(2))
    add_textbox(s6, Inches(1.8), y + Inches(0.15), Inches(3), Inches(0.4),
        name, Pt(18), TEXT_DARK, True)
    add_textbox(s6, Inches(5), y + Inches(0.15), Inches(6.5), Inches(1),
        desc, Pt(15), TEXT_GRAY, False)

add_textbox(s6, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '06', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 7: Chapter Divider - Mechanical
# ==========================================
s7 = add_blank_slide()
set_bg(s7, WHITE)
add_line(s7, Inches(0), Inches(0), Inches(0), Inches(7.5), ACCENT_BLUE_2, Pt(6))
# Car image (if available)
if os.path.exists(CAR_MODEL):
    add_img(s7, CAR_MODEL, Inches(5), Inches(1.8), Inches(4.5))
add_textbox(s7, Inches(3), Inches(2.8), Inches(8), Inches(0.8),
    '机械结构设计', Pt(36), TEXT_DARK, True, PP_ALIGN.CENTER)
add_line(s7, Inches(5.5), Inches(3.7), Inches(1.2), Pt(3), ACCENT_BLUE_2)
add_textbox(s7, Inches(3), Inches(4), Inches(8), Inches(0.5),
    '麦克纳姆轮底盘 · 自动接驳装置', Pt(16), TEXT_GRAY, False, PP_ALIGN.CENTER)

# ==========================================
# SLIDE 8: Mechanical - Chassis
# ==========================================
s8 = add_blank_slide()
set_bg(s8, WHITE)
add_textbox(s8, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '麦克纳姆轮底盘设计', Pt(30), TEXT_DARK, True)
add_line(s8, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

# Key features
features = [
    '全向移动：前后左右+斜向，无需转弯半径',
    '四轮独立驱动：四个直流无刷电机各控一轮',
    '轻量化框架：减重同时保证结构强度',
    '适合矿区狭窄道路和复杂地面条件',
]
for i, feat in enumerate(features):
    add_textbox(s8, Inches(0.8), Inches(1.5) + Inches(i * 0.5), Inches(6.5), Inches(0.45),
        f'· {feat}', Pt(16), TEXT_GRAY, False)

# Car chassis image
if os.path.exists(CAR_CHASSIS):
    add_img(s8, CAR_CHASSIS, Inches(7.8), Inches(1.5), Inches(4.5))

add_textbox(s8, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '08', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 9: Mechanical - Docking Device
# ==========================================
s9 = add_blank_slide()
set_bg(s9, WHITE)
add_textbox(s9, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '自动接驳装置——全文核心创新点', Pt(30), TEXT_DARK, True)
add_line(s9, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

add_textbox(s9, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5),
    '一个装置，同时解决三个问题', Pt(18), TEXT_GRAY, False)

docking_features = [
    ('缓冲吸能', '弹簧/橡胶阻尼结构\n吸收降落冲击力', ACCENT_BLUE),
    ('侧向纠偏', '可摆动托板 ±30mm\n补偿定位偏差', GREEN_ACCENT),
    ('自锁固定', '机械卡扣自锁\n运输中货物不滑落', RED_ACCENT),
]
for i, (title, desc, color) in enumerate(docking_features):
    x = Inches(0.8) + Inches(i * 4.1)
    add_rect(s9, x, Inches(2.1), Inches(3.7), Inches(2.5), CARD_BG, color, Pt(2))
    add_textbox(s9, x + Inches(0.3), Inches(2.3), Inches(3), Inches(0.4),
        title, Pt(22), TEXT_DARK, True, PP_ALIGN.CENTER)
    add_textbox(s9, x + Inches(0.3), Inches(2.9), Inches(3), Inches(1.2),
        desc, Pt(15), TEXT_GRAY, False, PP_ALIGN.CENTER)

# Docking image
if os.path.exists(CAR_DOCK):
    add_img(s9, CAR_DOCK, Inches(7.8), Inches(5), Inches(4.5), Inches(2))

add_textbox(s9, Inches(0.8), Inches(5.2), Inches(6), Inches(0.8),
    '设计工具：SolidWorks 三维建模，明确装配关系与关键安装位置', Pt(14), TEXT_GRAY, False)

add_textbox(s9, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '09', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 10: Chapter Divider - Control
# ==========================================
s10 = add_blank_slide()
set_bg(s10, WHITE)
add_line(s10, Inches(0), Inches(0), Inches(0), Inches(7.5), ACCENT_BLUE_2, Pt(6))
add_textbox(s10, Inches(3), Inches(2.8), Inches(8), Inches(0.8),
    '电控与传感器系统', Pt(36), TEXT_DARK, True, PP_ALIGN.CENTER)
add_line(s10, Inches(5.5), Inches(3.7), Inches(1.2), Pt(3), ACCENT_BLUE_2)
add_textbox(s10, Inches(3), Inches(4), Inches(8), Inches(0.5),
    '工控机 + PLC 双级控制架构 · 安全互锁与传感器布局', Pt(16), TEXT_GRAY, False, PP_ALIGN.CENTER)

# ==========================================
# SLIDE 11: Control System Architecture
# ==========================================
s11 = add_blank_slide()
set_bg(s11, WHITE)
add_textbox(s11, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '电控系统架构', Pt(30), TEXT_DARK, True)
add_line(s11, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

# Architecture flow
add_rect(s11, Inches(3), Inches(1.8), Inches(7), Inches(0.8), DARK_BG)
add_textbox(s11, Inches(3.5), Inches(1.9), Inches(6), Inches(0.6),
    '工控机（上位机 · ROS 导航算法）', Pt(18), WHITE, True, PP_ALIGN.CENTER)

add_textbox(s11, Inches(6), Inches(2.8), Inches(1.5), Inches(0.5),
    '↕ TCP/Modbus', Pt(14), ACCENT_BLUE, False, PP_ALIGN.CENTER)

add_rect(s11, Inches(3), Inches(3.3), Inches(7), Inches(0.8), DARKER_BG)
add_textbox(s11, Inches(3.5), Inches(3.4), Inches(6), Inches(0.6),
    'PLC（下位机 · 逻辑控制 · 安全互锁）', Pt(18), WHITE, True, PP_ALIGN.CENTER)

# Output lines
outputs = [('电机驱动', Inches(1.5)), ('传感器 & 继电器', Inches(5.5)), ('电磁锁 & 指示灯', Inches(9.2))]
for name, x in outputs:
    add_textbox(s11, x, Inches(4.4), Inches(2.5), Inches(0.4),
        name, Pt(14), TEXT_GRAY, True, PP_ALIGN.CENTER)
    # Arrow line
    add_line(s11, Inches(6.5), Inches(4.0), Inches(0), Inches(0.3), ACCENT_BLUE_2, Pt(1))

# I/O info
add_textbox(s11, Inches(0.8), Inches(5.3), Inches(12), Inches(1),
    '· 传感器配置：2D激光雷达 + IMU + 轮式里程计 + 超声波传感器\n· I/O规划：数字量输入×16 / 数字量输出×12 / 模拟量输入×4\n· 电源系统：48V动力电池 → DC-DC → 24V（控制）/ 12V（传感器）/ 5V（工控机）', Pt(14), TEXT_GRAY, False)

add_textbox(s11, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '11', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 12: Docking Sequence Logic
# ==========================================
s12 = add_blank_slide()
set_bg(s12, WHITE)
add_textbox(s12, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '接驳作业时序逻辑', Pt(30), TEXT_DARK, True)
add_line(s12, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

steps = ['① 到达\n指定点', '② 粗定位\nGPS引导', '③ 微调\n激光对准', '④ 接货台\n升起', '⑤ 锁紧\n确认', '⑥ 运输\n至卸货点', '⑦ 复位\n待命']
for i, step in enumerate(steps):
    x = Inches(0.5) + Inches(i * 1.85)
    add_rect(s12, x, Inches(2.5), Inches(1.6), Inches(1.8), CARD_BG, ACCENT_BLUE if i in [3,4] else LIGHT_GRAY, Pt(1.5))
    add_textbox(s12, x + Inches(0.1), Inches(2.7), Inches(1.4), Inches(1.2),
        step, Pt(14), TEXT_DARK if i in [3,4] else TEXT_GRAY, True, PP_ALIGN.CENTER)
    if i < 6:
        add_textbox(s12, x + Inches(1.6), Inches(3.2), Inches(0.25), Inches(0.3),
            '→', Pt(16), ACCENT_BLUE_2, True, PP_ALIGN.CENTER)

add_textbox(s12, Inches(0.8), Inches(5), Inches(12), Inches(0.8),
    'PLC 梯形图实现顺序控制逻辑，每步依次触发，上一步未完成则下一步不启动', Pt(15), TEXT_GRAY, False)

add_textbox(s12, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '12', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 13: Kinematics
# ==========================================
s13 = add_blank_slide()
set_bg(s13, WHITE)
add_textbox(s13, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '麦克纳姆轮运动学模型', Pt(30), TEXT_DARK, True)
add_line(s13, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

add_textbox(s13, Inches(0.8), Inches(1.4), Inches(6), Inches(0.5),
    '核心问题：四个轮子转速不一样才能走直线', Pt(18), TEXT_GRAY, False)

# Formula box
add_rect(s13, Inches(0.8), Inches(2.2), Inches(5.5), Inches(2), CARD_BG, ACCENT_BLUE, Pt(1))
add_textbox(s13, Inches(1.2), Inches(2.4), Inches(4.8), Inches(0.4),
    '逆运动学公式（给定车速 → 各轮转速）', Pt(14), TEXT_DARK, True)
add_textbox(s13, Inches(1.2), Inches(2.9), Inches(4.8), Inches(1.2),
    'ωᵢ = (vₓ ± vᵧ ± ωz · L) / R\n\n其中：vₓ,vᵧ 为期望速度分量，ωz 为旋转角速度\nL 为轮距参数，R 为轮半径，i=1,2,3,4', Pt(15), TEXT_GRAY, False)

# Speed matching
add_textbox(s13, Inches(7.2), Inches(2.2), Inches(5.5), Inches(2.5),
    '速度匹配关系推导结果：\n· 直行：四个轮等速同向\n· 横移：对角轮等速反向\n· 原地旋转：同侧轮反向\n· 斜行：内外轮速差补偿\n\n模型用于后续 PID 控制的目标值设定', Pt(15), TEXT_GRAY, False)

add_textbox(s13, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '13', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 14: PID Control
# ==========================================
s14 = add_blank_slide()
set_bg(s14, WHITE)
add_textbox(s14, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    'PID 转速闭环控制', Pt(30), TEXT_DARK, True)
add_line(s14, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

# Control diagram
add_rect(s14, Inches(0.8), Inches(1.6), Inches(8), Inches(1.8), CARD_BG, None)
add_textbox(s14, Inches(1.2), Inches(1.7), Inches(7.2), Inches(1.6),
    '目标转速  →  [前馈补偿（理论值）]  →  [PID控制器]  →  [直流电机模型]  →  实际转速\n                           ↑                                              ↓\n                           └────────── 转速反馈（编码器） ──────────┘',
    Pt(13), TEXT_GRAY, False)

# Key results
add_textbox(s14, Inches(0.8), Inches(3.8), Inches(5.5), Inches(0.5),
    '仿真结果（MATLAB）', Pt(18), TEXT_DARK, True)

# Result cards
results = [
    ('< 0.1s', '上升时间'),
    ('< 5%', '超调量'),
    ('< 2%', '稳态误差'),
    ('~40%', '响应时间减少\n（vs纯PID）'),
]
for i, (value, label) in enumerate(results):
    x = Inches(0.8) + Inches(i * 3.2)
    add_rect(s14, x, Inches(4.4), Inches(2.8), Inches(1.5), CARD_BG, None)
    add_textbox(s14, x + Inches(0.3), Inches(4.5), Inches(2.2), Inches(0.6),
        value, Pt(28), ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s14, x + Inches(0.3), Inches(5.1), Inches(2.2), Inches(0.5),
        label, Pt(13), TEXT_GRAY, False, PP_ALIGN.CENTER)

add_textbox(s14, Inches(0.8), Inches(6.2), Inches(10), Inches(0.6),
    '结论：前馈补偿+PID双层控制策略有效——响应更快，超调可控，稳态精度满足工程要求', Pt(15), TEXT_GRAY, False)

add_textbox(s14, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '14', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 15: Simulation Results
# ==========================================
s15 = add_blank_slide()
set_bg(s15, WHITE)
add_textbox(s15, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '仿真结果验证——PID 与路径规划', Pt(28), TEXT_DARK, True)
add_line(s15, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

# Top row: PID curves (placeholders for thesis images)
add_textbox(s15, Inches(0.8), Inches(1.5), Inches(5.5), Inches(0.4),
    'PID 阶跃响应曲线', Pt(16), TEXT_DARK, True)
add_rect(s15, Inches(0.8), Inches(2), Inches(5.5), Inches(2), CARD_BG, LIGHT_GRAY, Pt(1))
add_textbox(s15, Inches(2.5), Inches(2.8), Inches(2.5), Inches(0.4),
    '[论文图5-5：阶跃响应]', Pt(14), FAINT, False, PP_ALIGN.CENTER)

add_textbox(s15, Inches(7), Inches(1.5), Inches(5.5), Inches(0.4),
    'PID 跟踪响应曲线', Pt(16), TEXT_DARK, True)
add_rect(s15, Inches(7), Inches(2), Inches(5.5), Inches(2), CARD_BG, LIGHT_GRAY, Pt(1))
add_textbox(s15, Inches(8.7), Inches(2.8), Inches(2.5), Inches(0.4),
    '[论文图5-6：跟踪响应]', Pt(14), FAINT, False, PP_ALIGN.CENTER)

# Bottom row: A* comparison
add_textbox(s15, Inches(0.8), Inches(4.4), Inches(5.5), Inches(0.4),
    '传统 A* 路径', Pt(16), TEXT_DARK, True)
add_rect(s15, Inches(0.8), Inches(4.9), Inches(5.5), Inches(2), CARD_BG, LIGHT_GRAY, Pt(1))
add_textbox(s15, Inches(2.5), Inches(5.7), Inches(2.5), Inches(0.4),
    '[论文图6-5：传统A*]', Pt(14), FAINT, False, PP_ALIGN.CENTER)

add_textbox(s15, Inches(7), Inches(4.4), Inches(5.5), Inches(0.4),
    '改进 A* 平滑路径（B样条优化）', Pt(16), TEXT_DARK, True)
add_rect(s15, Inches(7), Inches(4.9), Inches(5.5), Inches(2), CARD_BG, LIGHT_GRAY, Pt(1))
add_textbox(s15, Inches(8.7), Inches(5.7), Inches(2.5), Inches(0.4),
    '[论文图6-6：改进A*]', Pt(14), FAINT, False, PP_ALIGN.CENTER)

add_textbox(s15, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '15', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 16: SLAM & Navigation
# ==========================================
s16 = add_blank_slide()
set_bg(s16, WHITE)
add_textbox(s16, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    'SLAM 建图与导航系统', Pt(30), TEXT_DARK, True)
add_line(s16, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

add_textbox(s16, Inches(0.8), Inches(1.4), Inches(10), Inches(0.5),
    '多传感器融合定位 + 分层路径规划', Pt(18), TEXT_GRAY, False)

# Three modules
nav_modules = [
    ('SLAM 建图', '· 2D激光雷达扫描环境\n· IMU提供姿态角\n· 轮式里程计推算位移\n· Cartographer算法融合\n→ 输出：二维栅格地图', ACCENT_BLUE),
    ('A* 全局路径', '· 栅格地图建模\n· 启发式搜索最短路径\n· B样条平滑处理\n· 去除冗余节点\n→ 输出：平滑全局路径', GREEN_ACCENT),
    ('DWA 局部避障', '· 实时速度空间采样\n· 多轨迹评分\n· 动态障碍物规避\n· 对接停靠精度控制\n→ 输出：实时速度指令', RED_ACCENT),
]
for i, (title, desc, color) in enumerate(nav_modules):
    x = Inches(0.8) + Inches(i * 4.2)
    add_rect(s16, x, Inches(2.3), Inches(3.8), Inches(4), CARD_BG, color, Pt(2))
    add_textbox(s16, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(0.4),
        title, Pt(20), TEXT_DARK, True)
    add_textbox(s16, x + Inches(0.3), Inches(3.1), Inches(3.2), Inches(2.8),
        desc, Pt(15), TEXT_GRAY, False)

add_textbox(s16, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '16', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 17: Innovation Points
# ==========================================
s17 = add_blank_slide()
set_bg(s17, WHITE)
add_textbox(s17, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '创新点', Pt(30), TEXT_DARK, True)
add_line(s17, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

innovations = [
    ('01', '缓冲纠偏一体化\n接驳装置', '同时解决降落冲击、定位偏差和货物固定三个问题，\n不是分开处理，而是单一装置集成', ACCENT_BLUE),
    ('02', '前馈 + PID\n双层调速策略', '前馈计算理论值直接输出，PID只修正偏差，\n响应时间比纯PID减少约40%', GREEN_ACCENT),
    ('03', 'A* 平滑 + DWA\n分层导航架构', '全局路径最优 + 局部实时避障，\nB样条平滑去除折线拐点', RED_ACCENT),
]
for i, (num, title, desc, color) in enumerate(innovations):
    y = Inches(1.8) + Inches(i * 1.75)
    add_rect(s17, Inches(1.5), y, Inches(10.5), Inches(1.5), CARD_BG, color, Pt(2))
    add_textbox(s17, Inches(1.8), y + Inches(0.2), Inches(0.8), Inches(0.6),
        num, Pt(32), color, True)
    add_textbox(s17, Inches(2.8), y + Inches(0.15), Inches(3.5), Inches(1.2),
        title, Pt(18), TEXT_DARK, True)
    add_textbox(s17, Inches(6.5), y + Inches(0.3), Inches(5.2), Inches(0.9),
        desc, Pt(14), TEXT_GRAY, False)

add_textbox(s17, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '17', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 18: Summary
# ==========================================
s18 = add_blank_slide()
set_bg(s18, WHITE)
add_textbox(s18, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '研究总结', Pt(30), TEXT_DARK, True)
add_line(s18, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

accomplishments = [
    '完成整车 SolidWorks 三维建模，装配关系验证通过',
    '设计缓冲纠偏一体化接驳装置（核心创新）',
    '搭建 PLC + 工控机电控架构，I/O规划与时序逻辑完成',
    'PID 仿真验证：超调量 < 5%，响应时间 < 0.1s',
    'A* + DWA 路径规划仿真通过，路径平滑有效',
]
for i, acc in enumerate(accomplishments):
    y = Inches(1.6) + Inches(i * 0.65)
    # Checkmark
    add_textbox(s18, Inches(1.5), y, Inches(0.4), Inches(0.4),
        '✓', Pt(18), GREEN_ACCENT, True)
    add_textbox(s18, Inches(2.1), y, Inches(9), Inches(0.4),
        acc, Pt(16), TEXT_GRAY, False)

# Key stats
add_line(s18, Inches(1.5), Inches(5.2), Inches(10), Pt(1), LIGHT_GRAY)
stats = [('52页', '论文篇幅'), ('28篇', '参考文献'), ('3项', '创新点'), ('SolidWorks\n+MATLAB', '工具链')]
for i, (value, label) in enumerate(stats):
    x = Inches(1.5) + Inches(i * 2.8)
    add_textbox(s18, x, Inches(5.5), Inches(2.2), Inches(0.5),
        value, Pt(22), ACCENT_BLUE, True, PP_ALIGN.CENTER)
    add_textbox(s18, x, Inches(5.9), Inches(2.2), Inches(0.4),
        label, Pt(13), TEXT_GRAY, False, PP_ALIGN.CENTER)

add_textbox(s18, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '18', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 19: Limitations & Future
# ==========================================
s19 = add_blank_slide()
set_bg(s19, WHITE)
add_textbox(s19, Inches(0.8), Inches(0.5), Inches(10), Inches(0.6),
    '不足与展望', Pt(30), TEXT_DARK, True)
add_line(s19, Inches(0.8), Inches(1.05), Inches(0.6), Pt(3), ACCENT_BLUE_2)

limitations = [
    ('只做仿真\n未造实物', '条件允许时制作1:5缩比样机\n进行实物测试验证'),
    ('矿区工况\n仅考虑平地', '后续加入坡度识别\n和功率补偿算法'),
    ('SLAM仅在\n室内环境测试', '实际露天矿区环境\n进行现场建图验证'),
]
for i, (limit, future) in enumerate(limitations):
    y = Inches(1.8) + Inches(i * 1.7)
    add_rect(s19, Inches(1.5), y, Inches(10.5), Inches(1.45), CARD_BG, None)
    add_textbox(s19, Inches(1.8), y + Inches(0.2), Inches(2.5), Inches(1),
        limit, Pt(16), RED_ACCENT, True)
    add_textbox(s19, Inches(4.5), y + Inches(0.15), Inches(1), Inches(0.4),
        '→', Pt(24), ACCENT_BLUE_2, True, PP_ALIGN.CENTER)
    add_textbox(s19, Inches(5.5), y + Inches(0.3), Inches(6), Inches(0.8),
        future, Pt(15), TEXT_GRAY, False)

add_textbox(s19, Inches(12.2), Inches(7.1), Inches(0.8), Inches(0.3),
    '19', Pt(10), FAINT, False, PP_ALIGN.RIGHT)

# ==========================================
# SLIDE 20: Thank You
# ==========================================
s20 = add_blank_slide()
set_bg(s20, DARK_BG)
add_line(s20, Inches(0.8), Inches(0.9), Inches(1.5), Pt(4), ACCENT_BLUE_2)
add_textbox(s20, Inches(3), Inches(2.5), Inches(7.5), Inches(1),
    '感谢各位老师，请批评指正', Pt(34), WHITE, True, PP_ALIGN.CENTER)
add_line(s20, Inches(5.5), Inches(3.8), Inches(2), Pt(2), ACCENT_BLUE_2)
add_multi_text(s20, Inches(4), Inches(4.3), Inches(5.5), Inches(1), [
    ('张紫然  ·  机器人22-1  ·  齐鲁工业大学机械工程学院', Pt(14), RGBColor(0xCC,0xCC,0xCC), False, PP_ALIGN.CENTER),
])
add_textbox(s20, Inches(5), Inches(6.6), Inches(3.5), Inches(0.4),
    '2026年5月', Pt(11), RGBColor(0x99,0x99,0x99), False, PP_ALIGN.CENTER)
add_line(s20, Inches(0.8), Inches(7.2), Inches(11.5), Pt(1), RGBColor(0x33,0x44,0x55))

# ==========================================
# SAVE
# ==========================================
output_path = 'D:/素材/Claude/portfolio/无人接驳小车-答辩PPT.pptx'
prs.save(output_path)
print(f'PPT saved: {output_path}')
print(f'Total slides: {len(prs.slides)}')
